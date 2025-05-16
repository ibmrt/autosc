import pyautogui
import threading
import keyboard
from datetime import datetime
from winotify import Notification
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from autosc_events import *


class Process:
    def __init__(self):
        self.key = True

    def terminate(self):
        self.key = False

    def __bool__(self) -> bool:
        return self.key

    def __repr__(self) -> str:
        return f"Process({self.key})"

class Engine:

    temp_coords = []

    @staticmethod
    def notify(msg):
        thread = threading.Thread(target=Notification("AutoSC", "提示", msg=msg, icon=r"C:\Users\xiqix\Desktop\AutoSC\jie.png",
                     duration="short").show(), args=())
        thread.daemon = True
        thread.start()

    @staticmethod
    def list_of_image(path:Path):
        gen = path.iterdir()
        list_of_image = [x for x in gen if str(x.suffix) in (".png", ".jpg", ".jpeg")]
        return list_of_image

    @staticmethod
    def convert_to_pixel(width, height) -> tuple[int, int]:
        dpi = 72
        slide_width_in = width / 914400
        slide_height_in = height / 914400

        slide_width_px = int(slide_width_in * dpi)
        slide_height_px = int(slide_height_in * dpi)
        return slide_width_px, slide_height_px

    @staticmethod
    def get_resize_dim(img_width, img_height, slide_width, slide_height):
        img_ratio = img_width / img_height
        slide_ratio = slide_width / slide_height
        if img_ratio < slide_ratio:
            multiplier = slide_height / img_height
            return int(img_width * multiplier), int(img_height * multiplier)
        else:
            multiplier = slide_width / img_width
            return int(img_width * multiplier), int(img_height * multiplier)

    @staticmethod
    def add_img(slide, img, img_width, img_height, slide_width, slide_height):
        if img_height == slide_height:
            slide.shapes.add_picture(img, Inches((slide_width / 72 - img_width / 72) / 2), 0, Inches(img_width / 72),
                                     Inches(img_height / 72))
            print((slide_width / 72 - img_width / 72) / 2)
        elif img_width == slide_width:
            slide.shapes.add_picture(img, 0, Inches((slide_height / 72 - img_height / 72) / 2), Inches(img_width / 72),
                                    Inches(img_height / 72))
        else:
            slide.shapes.add_picture(img, 0, 0, Inches(img_width / 72),
                                     Inches(img_height / 72))


    def __init__(self):
        self.box = []
        self.images = []
        self.save_path = None
        self.session = None

    def load_setting(self) -> None:
        settings = Path("settings.txt")
        if settings.exists():
            print("查询到已存在的设置文件：settings.txt")
            try:
                with open(settings, encoding="utf-8") as setting:
                    coords_list = setting.readline().split(", ")
                    coords = [int(coord.strip()) for coord in coords_list]
                    self.box = [(coords[i*2], coords[i*2+1]) for i in range(2)]
                    self.save_path = setting.readline()
                print(f"已载入既定设置:\n方框定位：{self.box}\n保存地点：{self.save_path}")
            except PermissionError:
                print("settings.txt存在——但没有阅读权限")
            
    def save_setting(self) -> SettingSaved|SaveSettingFailed:
        if self.box == [] or self.save_path is None:
            return SaveSettingFailed("信息不全：无法保存")
        try:
            coordinates = ", ".join([str(coord[i]) for coord in self.box for i in range(2)])
            with open("settings.txt", "w", encoding="utf-8") as setting:
                setting.write(f"{coordinates}\n{str(self.save_path)}")
                settings_txt = Path("settings.txt")
                Engine.notify(f"已把设置保存至{settings_txt.absolute()}")
                return SettingSaved(self.save_path, self.box)
        except PermissionError:
            return SaveSettingFailed("没有当前文件夹的编辑权限")
        except Exception as error:
            return SaveSettingFailed(error.__repr__())

    
    def get_path(self, path) -> None | SavePathFailed | PathSaved:
        if path == "":
            return SavePathFailed("输入路径空白")
        try:
            dest = Path(path)
            if dest.exists() and dest.is_dir():
                self.save_path = dest
                return PathSaved(str(dest))
            elif not dest.exists():
                return SavePathFailed("路径终点不存在")
            elif not dest.is_dir():
                return SavePathFailed("路径终点并非文件夹")
        except PermissionError:
            return SavePathFailed("程序没有该文件夹的阅读权限")


    def get_coords(self) -> SaveWindowSizeFailed | WindowSizeResult:
        Engine.temp_coords = []

        process = Process()
        keyboard.add_hotkey("ctrl+s", lambda: Engine.on_click(process))
        while process:
            pass
        keyboard.clear_hotkey("ctrl+s")
        if Engine.temp_coords[1][1]<Engine.temp_coords[0][1] or Engine.temp_coords[1][0]<Engine.temp_coords[0][0]:
            return SaveWindowSizeFailed("取景框的右下角与左上角错位")

        height = Engine.temp_coords[1][1] - Engine.temp_coords[0][1]
        width = Engine.temp_coords[1][0] - Engine.temp_coords[0][0]
        Engine.temp_coords = [(Engine.temp_coords[0][0], Engine.temp_coords[0][1]), (width, height)]
        self.box = Engine.temp_coords

        return WindowSizeResult(self.box[0][0], self.box[0][1], self.box[1][0], self.box[1][1])

    @staticmethod
    def on_click(process):
        x, y = pyautogui.position()
        Engine.temp_coords.append((x, y))
        Engine.notify(f"{x}, {y}")

        print(x, y)
        if len(Engine.temp_coords)==2:
            process.terminate()
    
    def flash(self, directory):
        num = 1
        while True:
            now = datetime.now()
            img = pyautogui.screenshot(region=(self.box[0][0], self.box[0][1],
                                               self.box[1][0], self.box[1][1]))
            if img not in self.images:
                self.images.append(img)
                self.images = self.images[-3:]
                img.save(directory/Path(fr"image_{num}.jpg"))
                yield f"@{num}: {now}"
                num += 1
            yield now
            pyautogui.sleep(1.5)


    @staticmethod
    def check_updates(process, gen):
        while process:
            print(next(gen))
        print("Terminated")
        exit()


    def run(self, process):
        self.session = None
        if self.save_path is None or self.box == []:
            return FailedToStartProcess("缺失路径或取景框坐标")
        now = datetime.now()
        dir_name = f"{now.year}_{now.month}_{now.day}_{now.hour}_{now.minute}"
        directory = Path(fr"{self.save_path}\{dir_name}")
        try:
            if directory.exists():
                directory.rmdir()
            directory.mkdir()
        except OSError:
            return FailedToStartProcess("检测到已存在的同名文件夹")
        self.session = self.flash(directory)
        thread = threading.Thread(target=Engine.check_updates, args=(process, self.session))
        thread.start()
        return StartedProcess(directory)

    def current_setting(self):
        return f"{self.save_path if self.save_path is not None else "暂无"}\n{self.box if self.box!=[] else "暂无"}"

    @staticmethod
    def make_ppt(path):
        dir_path = Path(path)
        if not dir_path.exists():
            return FailedToCreatePPT("未查询到该文件夹（等待系统保存或检查输入的路径）")
        elif not dir_path.is_dir():
            return FailedToCreatePPT("路径终点并非文件夹")
        else:
            image_path_list = Engine.list_of_image(dir_path)
            if not image_path_list:
                return FailedToCreatePPT("文件夹内没有图片")
            else:
                try:
                    prs = Presentation()
                    width, height = Engine.convert_to_pixel(prs.slide_width, prs.slide_height)
                    blank = prs.slide_layouts[6]
                    for img in image_path_list:
                        cur_slide = prs.slides.add_slide(blank)
                        image = Image.open(img)
                        img_width, img_height = Engine.get_resize_dim(image.width, image.height, width, height)
                        Engine.add_img(cur_slide, str(img), img_width, img_height, width, height)
                    save_to = str(dir_path/"presentation.pptx")
                    prs.save(save_to)
                    Engine.notify(f"保存至：{save_to}")
                    return save_to
                except PermissionError:
                    return FailedToCreatePPT("没有该文件夹的权限")
                except OSError as e:
                    return FailedToCreatePPT(e.__repr__())



    def process(self, event):
        if isinstance(event, StartAcquireWindow):
            return self.get_coords()
        elif isinstance(event, SaveFolderPath):
            return self.get_path(event.path)
        elif isinstance(event, SaveSettings):
            return self.save_setting()
        elif isinstance(event, StartProcess):
            return self.run(event.process)
        elif isinstance(event, StartCreatePPT):
            return self.make_ppt(event.path)